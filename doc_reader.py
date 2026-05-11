#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""读取文档/图片为纯文本，供 agent 使用。
用法:  python doc_reader.py <file>
输出:  JSON {ok, type, text, pages?, error?}
"""
import sys, os, json, traceback

def read_pdf(p):
    """提取 PDF 文本层 + 嵌入图片(可选 OCR) + 扫描页 OCR。
    优先 PyMuPDF (fitz)；缺失则退化到 pypdf。"""
    use_fitz = False
    try:
        import fitz  # PyMuPDF
        use_fitz = True
    except Exception:
        pass
    if not use_fitz:
        try:
            from pypdf import PdfReader
        except Exception:
            try:
                from PyPDF2 import PdfReader
            except Exception:
                raise RuntimeError("缺少 pypdf / PyMuPDF，请: pip install pymupdf 或 pip install pypdf")
        r = PdfReader(p)
        out = []
        for i, page in enumerate(r.pages):
            try:
                out.append(f"--- 第 {i+1} 页 ---\n" + (page.extract_text() or ""))
            except Exception as e:
                out.append(f"--- 第 {i+1} 页 (失败: {e}) ---")
        return {"type": "pdf", "pages": len(r.pages), "text": "\n\n".join(out),
                "image_count": 0, "ocr_used": False,
                "note": "未安装 PyMuPDF，无法提取图片/扫描页 OCR；如需: pip install pymupdf pytesseract"}

    # ---- PyMuPDF 路径 ----
    doc = fitz.open(p)
    n_pages = doc.page_count
    out = []
    n_images = 0
    image_brief = []
    ocr_pages = 0
    ocr_chars = 0
    ocr_error = None
    # OCR 引擎（可选）
    pytess = None
    try:
        import pytesseract
        from PIL import Image
        import io as _io
        pytess = (pytesseract, Image, _io)
    except Exception as e:
        ocr_error = "未安装 pytesseract/Pillow/Tesseract，跳过 OCR；安装: pip install pytesseract pillow + Tesseract OCR + 中文语言包"

    OCR_MAX_PAGES = int(os.environ.get('PDF_OCR_MAX_PAGES', '30'))
    OCR_MAX_CHARS = int(os.environ.get('PDF_OCR_MAX_CHARS', '200000'))
    EMBED_IMG_OCR = os.environ.get('PDF_EMBED_IMG_OCR', '1') != '0'

    # 图片导出目录（由 server.js 通过环境变量传入）
    IMG_OUT_DIR = os.environ.get('PDF_IMG_OUT_DIR') or ''
    SAVE_SCAN_PAGES = os.environ.get('PDF_SAVE_SCAN_PAGES', '1') != '0'
    if IMG_OUT_DIR:
        try:
            os.makedirs(IMG_OUT_DIR, exist_ok=True)
        except Exception:
            IMG_OUT_DIR = ''
    saved_images = []  # [{path, page, index, w, h, ext, kind, ocr?}]

    for i in range(n_pages):
        page = doc.load_page(i)
        try:
            txt = page.get_text("text") or ""
        except Exception as e:
            txt = f"(文本提取失败: {e})"
        page_out = [f"--- 第 {i+1} 页 ---", txt]

        # 列出嵌入图片
        try:
            imgs = page.get_images(full=True)
        except Exception:
            imgs = []
        if imgs:
            n_images += len(imgs)
            for j, img in enumerate(imgs):
                xref = img[0]
                try:
                    info = doc.extract_image(xref)
                    w = info.get('width'); h = info.get('height'); ext = info.get('ext') or 'png'
                    image_brief.append(f"[图 P{i+1}-{j+1} {w}x{h} .{ext}]")
                    saved_path = ''
                    if IMG_OUT_DIR:
                        try:
                            saved_path = os.path.join(IMG_OUT_DIR, f"p{i+1:03d}_{j+1:02d}.{ext}")
                            with open(saved_path, 'wb') as fp:
                                fp.write(info['image'])
                        except Exception:
                            saved_path = ''
                    img_ocr_text = ''
                    # 嵌入图 OCR（仅对足够大的位图，避免 logo/装饰小图）
                    if EMBED_IMG_OCR and pytess and ocr_chars < OCR_MAX_CHARS and w and h and w*h >= 200*200:
                        try:
                            pt, Image, _io = pytess
                            im = Image.open(_io.BytesIO(info['image']))
                            t = pt.image_to_string(im, lang="chi_sim+eng")
                            t = (t or '').strip()
                            if t:
                                page_out.append(f"[嵌入图 P{i+1}-{j+1} OCR]\n{t}")
                                ocr_chars += len(t)
                                img_ocr_text = t[:500]
                        except Exception as e:
                            if not ocr_error:
                                ocr_error = f"嵌入图 OCR 失败: {e}"
                    if saved_path:
                        saved_images.append({
                            'path': saved_path, 'page': i+1, 'index': j+1,
                            'w': w, 'h': h, 'ext': ext, 'kind': 'embedded',
                            'ocr': img_ocr_text,
                        })
                except Exception as e:
                    image_brief.append(f"[图 P{i+1}-{j+1} 提取失败: {e}]")

        # 文本层很少 → 视为扫描页，整页渲染 OCR
        if pytess and ocr_pages < OCR_MAX_PAGES and ocr_chars < OCR_MAX_CHARS and len(txt.strip()) < 40:
            try:
                pt, Image, _io = pytess
                pix = page.get_pixmap(dpi=200, alpha=False)
                im = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                t = pt.image_to_string(im, lang="chi_sim+eng")
                t = (t or '').strip()
                if t:
                    page_out.append(f"[扫描页 OCR P{i+1}]\n{t}")
                    ocr_chars += len(t); ocr_pages += 1
                # 同时保存扫描页栅格图（供前端预览）
                if IMG_OUT_DIR and SAVE_SCAN_PAGES:
                    try:
                        scan_path = os.path.join(IMG_OUT_DIR, f"scan_p{i+1:03d}.png")
                        pix.save(scan_path)
                        saved_images.append({
                            'path': scan_path, 'page': i+1, 'index': 0,
                            'w': pix.width, 'h': pix.height, 'ext': 'png',
                            'kind': 'scan_page', 'ocr': t[:500],
                        })
                    except Exception:
                        pass
            except Exception as e:
                if not ocr_error:
                    ocr_error = f"扫描页 OCR 失败: {e}"

        out.append("\n".join(page_out))

    text = "\n\n".join(out)
    if image_brief:
        text = f"[本 PDF 含 {n_images} 张嵌入图片]\n" + "\n".join(image_brief[:60]) + ("\n..." if len(image_brief) > 60 else "") + "\n\n" + text
    return {
        "type": "pdf", "pages": n_pages, "text": text,
        "image_count": n_images,
        "ocr_used": ocr_pages > 0 or (EMBED_IMG_OCR and pytess is not None and n_images > 0),
        "ocr_pages": ocr_pages,
        "ocr_chars": ocr_chars,
        "ocr_error": ocr_error,
        "images": saved_images,
        "image_dir": IMG_OUT_DIR or None,
    }

def read_docx(p):
    try:
        import docx  # python-docx
    except Exception:
        # 退化：解析 zip 里的 word/document.xml
        import zipfile, re, html
        with zipfile.ZipFile(p) as z:
            xml = z.read("word/document.xml").decode("utf-8", "ignore")
        text = re.sub(r"<[^>]+>", "", re.sub(r"</w:p>", "\n", xml))
        return {"type": "docx", "text": html.unescape(text)}
    d = docx.Document(p)
    paras = [pp.text for pp in d.paragraphs]
    for tb in d.tables:
        for row in tb.rows:
            paras.append(" | ".join(c.text for c in row.cells))
    return {"type": "docx", "text": "\n".join(paras)}

def read_doc(p):
    return {"type": "doc", "text": "(.doc 旧格式不支持，请另存为 .docx 或 PDF)"}

def read_pptx(p):
    try:
        from pptx import Presentation
    except Exception:
        raise RuntimeError("缺少 python-pptx，请: pip install python-pptx")
    prs = Presentation(p)
    out = []
    for i, slide in enumerate(prs.slides):
        out.append(f"--- 幻灯片 {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
    return {"type": "pptx", "text": "\n".join(out)}

def read_xlsx(p):
    try:
        from openpyxl import load_workbook
    except Exception:
        raise RuntimeError("缺少 openpyxl，请: pip install openpyxl")
    wb = load_workbook(p, data_only=True)
    out = []
    for name in wb.sheetnames:
        ws = wb[name]
        out.append(f"--- 工作表: {name} ---")
        for row in ws.iter_rows(values_only=True):
            out.append("\t".join("" if c is None else str(c) for c in row))
    return {"type": "xlsx", "text": "\n".join(out)}

def read_image(p):
    info = {"type": "image", "text": ""}
    try:
        from PIL import Image
        im = Image.open(p)
        info["size"] = list(im.size)
        info["mode"] = im.mode
    except Exception as e:
        info["pil_error"] = str(e)
    # OCR（可选）
    try:
        import pytesseract
        from PIL import Image as _I
        info["text"] = pytesseract.image_to_string(_I.open(p), lang="chi_sim+eng")
    except Exception as e:
        info["ocr_error"] = str(e)
        info["text"] = "(未安装 pytesseract / 中文 traineddata，已跳过 OCR；如需识别文字: pip install pytesseract，并安装 Tesseract OCR + 中文语言包)"
    return info

def read_text(p):
    for enc in ("utf-8", "gbk", "latin1"):
        try:
            with open(p, "r", encoding=enc) as f: return {"type": "text", "text": f.read()}
        except Exception: pass
    return {"type": "text", "text": "(无法以文本方式读取)"}

EXT_MAP = {
    ".pdf": read_pdf,
    ".docx": read_docx, ".doc": read_doc,
    ".pptx": read_pptx,
    ".xlsx": read_xlsx, ".xlsm": read_xlsx,
    ".png": read_image, ".jpg": read_image, ".jpeg": read_image,
    ".bmp": read_image, ".gif": read_image, ".webp": read_image, ".tif": read_image, ".tiff": read_image,
}

def main():
    if len(sys.argv) < 2:
        print(json.dumps({"ok": False, "error": "usage: doc_reader.py <file>"})); return
    p = sys.argv[1]
    if not os.path.exists(p):
        print(json.dumps({"ok": False, "error": f"文件不存在: {p}"})); return
    ext = os.path.splitext(p)[1].lower()
    fn = EXT_MAP.get(ext, read_text)
    try:
        out = fn(p); out["ok"] = True
        print(json.dumps(out, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({"ok": False, "error": str(e), "trace": traceback.format_exc()[-800:]}, ensure_ascii=False))

if __name__ == "__main__":
    main()
